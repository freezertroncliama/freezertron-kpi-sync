# -*- coding: utf-8 -*-
"""
Gera o Relatório Mensal PMOC (.pptx) no mesmo layout/design usado até
julho/2026 — o processo externo que gerava esse arquivo (fora dos dois
repos do projeto) parou de rodar depois disso, então este script assume
o lugar dele.

Como funciona: usa `modelo_relatorio_pmoc.pptx` (cópia de 2026-07.pptx,
o último relatório gerado pelo processo antigo) como MOLDE visual — não
recria nada do zero. Só troca o conteúdo (números, textos, linhas de
tabela, cards e o gráfico de tendência) mantendo cores, fontes e layout
originais intactos.

Fonte dos dados: as mesmas tabelas do Supabase que alimentam o painel
/kpi-manutencao (kpi_manutencao_atendimentos e kpi_manutencao_equipamentos)
— não mexe em Nextcloud/PDF diretamente, então só funciona pra um mês
depois que o kpi_manutencao.py já tiver rodado (sincronizado) esse mês.

Uso:
    python gerar_relatorio_pptx.py --cliente Mangels --mes 2026-08
    python gerar_relatorio_pptx.py --cliente Mangels          (usa o mês
        anterior ao atual por padrão)
    python gerar_relatorio_pptx.py --cliente Mangels --mes 2026-08 \
        --saida C:\\caminho\\2026-08.pptx
    python gerar_relatorio_pptx.py --cliente Mangels --upload  (além de
        gerar, sobe pro bucket relatorios-pmoc do Supabase — e pro
        Nextcloud, se as variáveis NEXTCLOUD_* estiverem configuradas —
        no nome AAAA-MM.pptx que o painel /kpi-manutencao já espera)

Requer no .env: SUPABASE_URL, SUPABASE_SERVICE_ROLE_KEY
Opcional (só se usar --upload e quiser cópia também no Nextcloud):
    NEXTCLOUD_URL, NEXTCLOUD_USER, NEXTCLOUD_APP_PASSWORD
"""
from __future__ import annotations

import argparse
import copy
import os
from collections import defaultdict
from datetime import datetime, timedelta

import requests
from dotenv import load_dotenv
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

load_dotenv()

SUPABASE_URL = os.environ["SUPABASE_URL"]
SUPABASE_SERVICE_ROLE_KEY = os.environ["SUPABASE_SERVICE_ROLE_KEY"]

DIR_ATUAL = os.path.dirname(os.path.abspath(__file__))
MODELO_PATH = os.path.join(DIR_ATUAL, "modelo_relatorio_pmoc.pptx")

CATEGORIA_CHILLER = "CHILLER\u00b4S"
PERIODICIDADE_DIAS = {"semanal": 7, "mensal": 30, "trimestral": 90}
# Mesmo limite visual do modelo original (slide de Pendências do Mês: AR
# mostrou 12 linhas + "+ 2 outro(s)" pra um total de 14).
MAX_ITENS_PENDENCIA_POR_CATEGORIA = 12
LIMIAR_REINCIDENCIA_DIAS = 30

MESES_PT = {
    1: "Janeiro", 2: "Fevereiro", 3: "Mar\u00e7o", 4: "Abril", 5: "Maio", 6: "Junho",
    7: "Julho", 8: "Agosto", 9: "Setembro", 10: "Outubro", 11: "Novembro", 12: "Dezembro",
}
MESES_PT_ABREV = {
    1: "jan", 2: "fev", 3: "mar", 4: "abr", 5: "mai", 6: "jun",
    7: "jul", 8: "ago", 9: "set", 10: "out", 11: "nov", 12: "dez",
}
NOMES_CLIENTE = {
    "Mangels": "Mangels S/A \u2014 Refrigera\u00e7\u00e3o A\u00e7os",
    "Noel": "Noel Supermercados",
}
NOMES_CLIENTE_CURTO = {
    "Mangels": "Mangels S/A",
    "Noel": "Noel Supermercados",
}


def nome_mes_ano(ano: int, mes: int) -> str:
    return f"{MESES_PT[mes]} {ano}"


def formatar_ptbr(v: float, casas: int = 0) -> str:
    s = f"{v:,.{casas}f}"
    s = s.replace(",", "_").replace(".", ",").replace("_", ".")
    return s


def formatar_horas(v: float) -> str:
    return formatar_ptbr(v, 1) + "h"


def fim_do_mes(ano: int, mes: int) -> datetime:
    if mes == 12:
        return datetime(ano, 12, 31)
    return datetime(ano, mes + 1, 1) - timedelta(days=1)


def parse_data(s: str | None) -> datetime | None:
    if not s:
        return None
    return datetime.fromisoformat(s[:10])


# ---------------------------------------------------------------------------
# 1) Dados — mesmas tabelas do Supabase que o painel /kpi-manutencao usa
# ---------------------------------------------------------------------------

def _get(tabela: str, params: dict) -> list[dict]:
    headers = {
        "apikey": SUPABASE_SERVICE_ROLE_KEY,
        "Authorization": f"Bearer {SUPABASE_SERVICE_ROLE_KEY}",
    }
    resp = requests.get(f"{SUPABASE_URL}/rest/v1/{tabela}", headers=headers, params=params)
    resp.raise_for_status()
    return resp.json()


def carregar_atendimentos(cliente: str) -> list[dict]:
    return _get(
        "kpi_manutencao_atendimentos",
        {
            "select": "categoria,equipamento,loja,arquivo,tipo,data,duracao_minutos,vazamento,gas_kg",
            "cliente": f"eq.{cliente}",
        },
    )


def carregar_equipamentos(cliente: str) -> list[dict]:
    return _get(
        "kpi_manutencao_equipamentos",
        {
            "select": "categoria,equipamento,loja,periodicidade,ultima_preventiva,"
            "ultima_substituicao_filtro,proxima_substituicao_filtro",
            "cliente": f"eq.{cliente}",
        },
    )


# ---------------------------------------------------------------------------
# 2) Cálculos — portados de app/kpi-manutencao/page.tsx (cobertura,
#    pendências, reincidência, tendência) + métricas novas específicas do
#    relatório (MTTR, ranking por horas corretivas, atendimentos no ano),
#    seguindo a mesma convenção usada no modelo de julho/2026.
# ---------------------------------------------------------------------------

CATEGORIAS_VALIDAS = {"AR CONDICIONADOS", "BEBEDOUROS", CATEGORIA_CHILLER}


def calcular_relatorio(cliente: str, atendimentos: list[dict], equipamentos: list[dict], ano: int, mes: int) -> dict:
    fim_mes = fim_do_mes(ano, mes)

    # Descarta entradas cuja "categoria"/"equipamento" não é uma das três
    # categorias reais do contrato — visto em produção: a pasta
    # "RELATÓRIOS PMOC" do Nextcloud (onde ficam os próprios relatórios
    # consolidados) é lida pela estrutura como se fosse um equipamento,
    # inflando "Equipamentos no Contrato" e furando a soma por categoria.
    equipamentos = [e for e in equipamentos if e["categoria"] in CATEGORIAS_VALIDAS]
    atendimentos = [a for a in atendimentos if a["categoria"] in CATEGORIAS_VALIDAS]

    for a in atendimentos:
        a["_data"] = parse_data(a.get("data"))

    validos = [a for a in atendimentos if a["tipo"] != "desconhecida"]
    do_mes = [a for a in validos if a["_data"] and a["_data"].year == ano and a["_data"].month == mes]
    do_ano_ate_mes = [a for a in validos if a["_data"] and a["_data"].year == ano and a["_data"] <= fim_mes]

    preventivas_mes = [a for a in do_mes if a["tipo"] == "preventiva"]
    corretivas_mes = sorted(
        (a for a in do_mes if a["tipo"] == "corretiva"), key=lambda a: a["_data"]
    )

    # --- agregados por equipamento (todo o histórico carregado) ---
    por_equip: dict[tuple[str, str], dict] = {}
    for a in validos:
        chave = (a["categoria"], a["equipamento"])
        if chave not in por_equip:
            por_equip[chave] = {
                "categoria": a["categoria"], "equipamento": a["equipamento"],
                "corretivas": 0, "horas": 0.0, "datas_corretiva": [],
            }
        info = por_equip[chave]
        if a["tipo"] == "corretiva":
            info["corretivas"] += 1
            info["horas"] += (a.get("duracao_minutos") or 0) / 60
            if a["_data"]:
                info["datas_corretiva"].append(a["_data"])

    candidatos = [v for v in por_equip.values() if v["corretivas"] > 0]
    candidatos.sort(key=lambda v: (v["horas"], v["corretivas"]), reverse=True)
    ranking = candidatos[:5]

    reincidencias = []
    for v in por_equip.values():
        datas = sorted(v["datas_corretiva"])
        if len(datas) < 2:
            continue
        menor = min((datas[i] - datas[i - 1]).days for i in range(1, len(datas)))
        if menor <= LIMIAR_REINCIDENCIA_DIAS:
            reincidencias.append({
                "categoria": v["categoria"], "equipamento": v["equipamento"],
                "total_corretivas": len(datas), "menor_intervalo": menor,
            })
    reincidencias.sort(key=lambda r: r["menor_intervalo"])

    # --- estrutura de equipamentos sob contrato ---
    estrutura: dict[tuple[str, str], dict] = {}
    for e in equipamentos:
        estrutura[(e["categoria"], e["equipamento"])] = e
    for a in validos:
        chave = (a["categoria"], a["equipamento"])
        estrutura.setdefault(chave, {
            "categoria": a["categoria"], "equipamento": a["equipamento"],
            "periodicidade": None, "ultima_preventiva": None,
        })

    preventivas_por_equip: dict[tuple[str, str], list[datetime]] = defaultdict(list)
    for a in validos:
        if a["tipo"] == "preventiva" and a["_data"]:
            preventivas_por_equip[(a["categoria"], a["equipamento"])].append(a["_data"])

    pendentes_por_categoria: dict[str, list[tuple[str, str]]] = defaultdict(list)
    total_contrato = len(estrutura)
    total_em_dia = 0
    for chave, e in estrutura.items():
        datas_prev = preventivas_por_equip.get(chave, [])
        teve_no_mes = any(d.year == ano and d.month == mes for d in datas_prev)
        if teve_no_mes:
            total_em_dia += 1
            continue
        ultima = max(datas_prev) if datas_prev else parse_data(e.get("ultima_preventiva"))
        if not ultima:
            detalhe = "Nunca teve preventiva"
        else:
            dias_periodicidade = PERIODICIDADE_DIAS.get(e.get("periodicidade"), 30)
            prazo = ultima + timedelta(days=dias_periodicidade)
            atraso = (fim_mes - prazo).days
            detalhe = f"{atraso}d atrasado" if atraso > 0 else "Sem preventiva no mes"
        pendentes_por_categoria[e["categoria"]].append((e["equipamento"], detalhe))

    for lista in pendentes_por_categoria.values():
        lista.sort(key=lambda t: t[0])

    total_pendente = total_contrato - total_em_dia
    pct_cumprimento = round(100 * total_em_dia / total_contrato) if total_contrato else 0

    cumprimento_por_categoria = []
    for categoria in sorted({e["categoria"] for e in estrutura.values()}):
        do_cat = [e for e in estrutura.values() if e["categoria"] == categoria]
        pendentes_cat = len(pendentes_por_categoria.get(categoria, []))
        realizado_cat = len(do_cat) - pendentes_cat
        cumprimento_por_categoria.append((categoria, realizado_cat, len(do_cat)))

    preventivas_por_categoria: dict[str, int] = defaultdict(int)
    for a in preventivas_mes:
        preventivas_por_categoria[a["categoria"]] += 1

    corretivas_com_duracao = [a for a in corretivas_mes if a.get("duracao_minutos")]
    if corretivas_com_duracao:
        mttr_horas = sum(a["duracao_minutos"] for a in corretivas_com_duracao) / 60 / len(corretivas_com_duracao)
        mttr_texto = formatar_horas(mttr_horas)
        mttr_sub = f"baseado em {len(corretivas_com_duracao)} corretiva(s) com registro"
    elif corretivas_mes:
        mttr_texto = "N/D"
        mttr_sub = f"{len(corretivas_mes)} corretiva(s) sem hor\u00e1rio de in\u00edcio/fim registrado"
    else:
        mttr_texto = "N/D"
        mttr_sub = "nenhuma corretiva no m\u00eas"

    hoje = datetime.now()
    bebedouros = [e for e in estrutura.values() if e["categoria"] == "BEBEDOUROS"]
    vencidos = []
    for e in bebedouros:
        prox = parse_data(e.get("proxima_substituicao_filtro"))
        if not prox:
            continue
        dias_para_vencer = (prox - hoje).days
        if dias_para_vencer < 0:
            vencidos.append({
                "equipamento": e["equipamento"], "proxima": prox, "dias_atrasado": -dias_para_vencer,
            })
    vencidos.sort(key=lambda v: -v["dias_atrasado"])

    por_mes: dict[str, dict] = defaultdict(lambda: {"preventivas": 0, "corretivas": 0})
    for a in validos:
        if not a["_data"] or a["_data"] > fim_mes:
            continue
        chave_mes = f"{a['_data'].year:04d}-{a['_data'].month:02d}"
        if a["tipo"] == "preventiva":
            por_mes[chave_mes]["preventivas"] += 1
        elif a["tipo"] == "corretiva":
            por_mes[chave_mes]["corretivas"] += 1
    meses_ordenados = sorted(por_mes.keys())[-4:]

    return {
        "cliente_nome": NOMES_CLIENTE.get(cliente, cliente),
        "cliente_nome_curto": NOMES_CLIENTE_CURTO.get(cliente, cliente),
        "nome_mes_ano": nome_mes_ano(ano, mes),
        "total_contrato": total_contrato,
        "total_em_dia": total_em_dia,
        "total_pendente": total_pendente,
        "pct_cumprimento": pct_cumprimento,
        "cumprimento_por_categoria": cumprimento_por_categoria,
        "pendentes_por_categoria": pendentes_por_categoria,
        "preventivas_mes_qtd": len(preventivas_mes),
        "corretivas_mes_qtd": len(corretivas_mes),
        "preventivas_por_categoria": preventivas_por_categoria,
        "corretivas_mes": corretivas_mes,
        "ranking": ranking,
        "reincidencias": reincidencias,
        "mttr_texto": mttr_texto,
        "mttr_sub": mttr_sub,
        "vencidos": vencidos,
        "atendimentos_no_ano": len(do_ano_ate_mes),
        "tendencia": [(m, por_mes[m]["preventivas"], por_mes[m]["corretivas"]) for m in meses_ordenados],
    }


# ---------------------------------------------------------------------------
# 3) Manipulação do PPTX — molde + substituição de conteúdo
# ---------------------------------------------------------------------------

def local_de_equipamento(equipamento: str) -> str:
    """'B015 - ACABAMENTO BOTIJÃO 01' -> 'Acabamento Botijão 01'
    (mesma transformação vista no modelo original: tira o código, capitaliza)."""
    nome = equipamento.split(" - ", 1)[-1] if " - " in equipamento else equipamento
    palavras = nome.split(" ")
    return " ".join(p.capitalize() if p else p for p in palavras)


def set_run_text(paragraph, texto: str) -> None:
    runs = paragraph.runs
    if not runs:
        run = paragraph.add_run()
    else:
        run = runs[0]
        for extra in runs[1:]:
            extra._r.getparent().remove(extra._r)
    run.text = texto


def set_shape_text(shape, texto: str, paragraph_index: int = 0) -> None:
    set_run_text(shape.text_frame.paragraphs[paragraph_index], texto)


def set_shape_paragraphs(shape, textos: list[str]) -> None:
    """Redimensiona um textbox de múltiplos parágrafos (clona/remove <a:p>
    usando o último parágrafo como molde de formatação) e escreve `textos`."""
    tf = shape.text_frame
    txBody = tf._txBody
    if not textos:
        for p in list(tf.paragraphs)[1:]:
            txBody.remove(p._p)
        set_run_text(tf.paragraphs[0], "")
        return
    template_p = tf.paragraphs[-1]._p
    while len(tf.paragraphs) < len(textos):
        txBody.append(copy.deepcopy(template_p))
    while len(tf.paragraphs) > len(textos):
        txBody.remove(tf.paragraphs[-1]._p)
    for p, texto in zip(tf.paragraphs, textos):
        set_run_text(p, texto)


def resize_table(table, n_linhas_dados: int) -> None:
    tbl = table._tbl
    trs = tbl.findall(qn("a:tr"))
    dados = trs[1:]  # trs[0] é o cabeçalho
    template = dados[-1]
    while len(dados) < n_linhas_dados:
        novo = copy.deepcopy(template)
        tbl.append(novo)
        dados.append(novo)
    while len(dados) > n_linhas_dados:
        tbl.remove(dados.pop())


def set_table_row(table, row_idx: int, valores: list[str]) -> None:
    row = table.rows[row_idx]
    for cell, valor in zip(row.cells, valores):
        set_run_text(cell.text_frame.paragraphs[0], valor)


def clonar_shapes(slide, indices: list[int], offset_y_emu: int):
    """Clona um conjunto de shapes (pelos índices atuais em slide.shapes) e
    desloca as cópias em Y — usado pros cards de filtro vencido, que
    aparecem em quantidade variável mês a mês."""
    shapes_atuais = list(slide.shapes)
    spTree = slide.shapes._spTree
    novos = []
    for idx in indices:
        el = copy.deepcopy(shapes_atuais[idx]._element)
        spTree.append(el)
        xfrm = el.find(f".//{qn('a:xfrm')}")
        if xfrm is not None:
            off = xfrm.find(qn("a:off"))
            off.set("y", str(int(off.get("y")) + offset_y_emu))
        novos.append(el)
    return novos


def mover_para_o_fim(slide, shape) -> None:
    slide.shapes._spTree.append(shape._element)


def texto_por_indice(slide, idx: int) -> str:
    return list(slide.shapes)[idx].text_frame.text


def atualizar_rodapes(prs: Presentation, texto_mes: str) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.startswith("Freezertron \u00b7 Relat\u00f3rio PMOC \u00b7"):
                set_shape_text(shape, f"Freezertron \u00b7 Relat\u00f3rio PMOC \u00b7 {texto_mes}")


def gerar_pptx(dados: dict, caminho_saida: str) -> None:
    prs = Presentation(MODELO_PATH)
    slides = list(prs.slides)
    d = dados

    # --- slide 0: capa (parágrafo composto por múltiplos runs) ---
    shape = list(slides[0].shapes)[5]
    runs = shape.text_frame.paragraphs[0].runs
    runs[1].text = f"{d['cliente_nome']}\n"
    runs[3].text = f"{d['total_contrato']} equipamentos\n"
    runs[5].text = d["nome_mes_ano"]

    # --- slide 1: apresentação (nome do cliente no texto de introdução) ---
    shape = list(slides[1].shapes)[3]
    texto_intro = shape.text_frame.text.replace("Mangels S/A", d["cliente_nome_curto"])
    set_shape_text(shape, texto_intro)

    # --- slide 2: riscos ---
    n_vencidos = len(d["vencidos"])
    if n_vencidos:
        max_dias = d["vencidos"][0]["dias_atrasado"]
        texto_risco = (
            f"{n_vencidos} bebedouro(s) com troca de filtro vencida h\u00e1 mais de "
            f"{max_dias} dias \u2014 ver detalhe no fim do relat\u00f3rio."
        )
    else:
        texto_risco = "Nenhum bebedouro com troca de filtro vencida nesta compet\u00eancia."
    set_shape_text(list(slides[2].shapes)[12], texto_risco)

    # --- slide 3: resumo executivo ---
    shapes = list(slides[3].shapes)
    set_shape_text(shapes[6], str(d["total_contrato"]))
    set_shape_text(shapes[14], str(d["atendimentos_no_ano"]))

    # --- slide 4: KPIs preventiva ---
    shapes = list(slides[4].shapes)
    set_shape_text(shapes[6], str(d["total_contrato"]))
    set_shape_text(shapes[9], str(d["total_em_dia"]))
    set_shape_text(shapes[12], f"{d['pct_cumprimento']}%")
    # (idx_texto, idx_barra_preenchida) \u2014 largura total da trilha (fundo) \u00e9
    # sempre 7223760 EMU no molde; a barra azul \u00e9 proporcional a realizado/total.
    LARGURA_TRILHA_SLIDE4 = 7223760
    indices_categoria = [(15, 17), (18, 20), (21, 23)]
    for (idx_texto, idx_barra), (categoria, realizado, total) in zip(indices_categoria, d["cumprimento_por_categoria"]):
        set_shape_text(shapes[idx_texto], f"{categoria} \u2014 {realizado}/{total}")
        frac = (realizado / total) if total else 0
        shapes[idx_barra].width = round(LARGURA_TRILHA_SLIDE4 * frac)

    # --- slide 5: KPIs corretiva ---
    shapes = list(slides[5].shapes)
    set_shape_text(shapes[6], str(d["corretivas_mes_qtd"]))
    set_shape_text(shapes[9], str(len(d["reincidencias"])))
    set_shape_text(shapes[13], d["mttr_texto"])
    set_shape_text(shapes[14], d["mttr_sub"])

    # --- slide 6: corretivas do mês ---
    shapes = list(slides[6].shapes)
    set_shape_text(
        shapes[1],
        f"{d['corretivas_mes_qtd']} chamado(s) corretivo(s) em {d['nome_mes_ano']} \u2014 aten\u00e7\u00e3o priorit\u00e1ria",
    )
    tabela = shapes[2].table
    linhas = d["corretivas_mes"]
    if linhas:
        resize_table(tabela, len(linhas))
        for i, a in enumerate(linhas):
            set_table_row(tabela, i + 1, [
                a["_data"].strftime("%d/%m"),
                a["categoria"],
                a["equipamento"],
                local_de_equipamento(a["equipamento"]),
            ])
    else:
        resize_table(tabela, 1)
        set_table_row(tabela, 1, ["\u2014", "\u2014", "Nenhuma corretiva no m\u00eas", "\u2014"])

    # --- slide 7: preventivas do mês ---
    shapes = list(slides[7].shapes)
    set_shape_text(
        shapes[1],
        f"{d['preventivas_mes_qtd']} manuten\u00e7\u00e3o(\u00f5es) preventiva(s) realizada(s) em "
        f"{d['nome_mes_ano']}, por categoria",
    )
    indices_cat_prev = {"AR CONDICIONADOS": 4, "BEBEDOUROS": 8, CATEGORIA_CHILLER: 12}
    for categoria, idx in indices_cat_prev.items():
        set_shape_text(shapes[idx], str(d["preventivas_por_categoria"].get(categoria, 0)))

    # --- slide 8: ranking (5 posições fixas) ---
    shapes = list(slides[8].shapes)
    # (nome, categoria, idx_barra, badge) \u2014 barra vermelha proporcional \u00e0s
    # horas corretivas relativas ao 1o colocado (largura total = 5486400
    # EMU no molde); mant\u00e9m uma largura m\u00ednima vis\u00edvel mesmo em 0h, igual
    # ao comportamento observado no relat\u00f3rio original.
    LARGURA_TRILHA_SLIDE8 = 5486400
    LARGURA_MINIMA_BARRA = 137160
    posicoes = [(3, 4, 6, 7), (9, 10, 12, 13), (15, 16, 18, 19), (21, 22, 24, 25), (27, 28, 30, 31)]
    ranking = d["ranking"]
    max_horas = max((item["horas"] for item in ranking), default=0)
    for (idx_nome, idx_cat, idx_barra, idx_badge), item in zip(posicoes, ranking + [None] * 5):
        if item is None:
            set_shape_text(shapes[idx_nome], "\u2014")
            set_shape_text(shapes[idx_cat], "\u2014")
            set_shape_text(shapes[idx_badge], "\u2014")
            shapes[idx_barra].width = LARGURA_MINIMA_BARRA
        else:
            set_shape_text(shapes[idx_nome], item["equipamento"])
            set_shape_text(shapes[idx_cat], item["categoria"])
            n_os = item["corretivas"]
            set_shape_text(shapes[idx_badge], f"{formatar_horas(item['horas'])} \u00b7 {n_os} OS")
            frac = (item["horas"] / max_horas) if max_horas else 0
            shapes[idx_barra].width = max(round(LARGURA_TRILHA_SLIDE8 * frac), LARGURA_MINIMA_BARRA)

    # --- slide 9: reincidência ---
    tabela = list(slides[9].shapes)[2].table
    reinc = d["reincidencias"]
    if reinc:
        resize_table(tabela, len(reinc))
        for i, r in enumerate(reinc):
            set_table_row(tabela, i + 1, [
                r["categoria"], r["equipamento"], str(r["total_corretivas"]), f"{r['menor_intervalo']}d",
            ])
    else:
        resize_table(tabela, 1)
        set_table_row(tabela, 1, ["\u2014", "Nenhuma reincid\u00eancia no per\u00edodo", "\u2014", "\u2014"])

    # --- slide 10: filtros de bebedouro vencidos (cards em quantidade variável) ---
    slide10 = slides[10]
    shapes = list(slide10.shapes)
    footer10 = shapes[12]
    OFFSET_CARD = 1005840  # top(card2) - top(card1), medido no molde
    TOPO_CARD1 = 1828800
    # Com 5+ cards o último encosta no rodapé (espaço vertical do slide é
    # fixo) — mostra só os N piores (já vêm ordenados por atraso desc.) e
    # um aviso "+ outros" pro resto, mesmo padrão do slide de Pendências.
    MAX_CARDS_FILTRO = 4
    vencidos_todos = d["vencidos"]
    vencidos = vencidos_todos[:MAX_CARDS_FILTRO]
    resto_filtro = len(vencidos_todos) - len(vencidos)
    if not vencidos:
        # some os dois cards de exemplo e mostra um aviso único
        for idx in (2, 3, 4, 5, 6, 7, 8, 9, 10, 11):
            el = shapes[idx]._element
            el.getparent().remove(el)
        mover_para_o_fim(slide10, footer10)
    else:
        indices_card1 = [2, 3, 4, 5, 6]
        n_extra = max(0, len(vencidos) - 2)
        for i in range(n_extra):
            clonar_shapes(slide10, indices_card1, OFFSET_CARD * (2 + i))
        # clonar_shapes anexa no fim da árvore de shapes — como o rodapé
        # ainda está "no meio" (era o último elemento ANTES da clonagem),
        # os clones acabam depois dele. Reposiciona pro fim agora, antes de
        # recalcular os índices, senão todo bloco clonado fica deslocado.
        mover_para_o_fim(slide10, footer10)
        shapes = list(slide10.shapes)
        # cards existentes: [2..6]=#1, [7..11]=#2, clones a partir do índice 12.
        # Ordem real de cada card (ver geometria do molde): [0]=fundo,
        # [1]=faixa colorida à esquerda, [2]=título, [3]=subtítulo, [4]=badge.
        blocos = [[2, 3, 4, 5, 6], [7, 8, 9, 10, 11]] + [
            [12 + 5 * i + j for j in range(5)] for i in range(n_extra)
        ]
        for bloco, v in zip(blocos, vencidos):
            _, _, idx_titulo, idx_sub, idx_badge = bloco
            set_shape_text(shapes[idx_titulo], v["equipamento"])
            set_shape_text(shapes[idx_sub], f"Pr\u00f3xima substitui\u00e7\u00e3o prevista: {v['proxima']:%d/%m/%Y}")
            set_shape_text(shapes[idx_badge], f"{v['dias_atrasado']}d atrasado")
        # esconde blocos de exemplo não usados (só acontece se vencidos == 1)
        for bloco in blocos[len(vencidos):]:
            for idx in bloco:
                el = shapes[idx]._element
                el.getparent().remove(el)
    if resto_filtro > 0:
        caixa = slide10.shapes.add_textbox(
            548640, TOPO_CARD1 + len(vencidos) * OFFSET_CARD, 11064240, 320040
        )
        # sem isso, o autofit padrão do add_textbox encolhe a caixa em volta
        # do texto e a recentraliza na posição original — parece "centralizado"
        # mesmo com alinhamento à esquerda.
        caixa.text_frame.word_wrap = True
        paragrafo = caixa.text_frame.paragraphs[0]
        paragrafo.alignment = PP_ALIGN.LEFT
        run = paragrafo.add_run()
        run.text = f"+ {resto_filtro} outro(s) bebedouro(s) com filtro vencido — lista completa em /kpi-manutencao."
        run.font.size = Pt(13)
        run.font.italic = True
        run.font.color.rgb = RGBColor(0x4A, 0x55, 0x68)
        run.font.name = "Calibri"
        mover_para_o_fim(slide10, footer10)

    # --- slide 11: tendência (gráfico nativo) ---
    chart = list(slides[11].shapes)[2].chart
    tend = d["tendencia"]
    dados_grafico = CategoryChartData()
    dados_grafico.categories = [
        f"{MESES_PT_ABREV[int(m.split('-')[1])].capitalize()}/{m.split('-')[0][2:]}" for m, _, _ in tend
    ]
    dados_grafico.add_series("Preventivas", [p for _, p, _ in tend])
    dados_grafico.add_series("Corretivas", [c for _, _, c in tend])
    chart.replace_data(dados_grafico)

    # --- slide 12: pendências do mês ---
    shapes = list(slides[12].shapes)
    set_shape_text(shapes[1], f"{d['total_pendente']} equipamento(s) ainda fora do prazo esperado de preventiva")
    blocos_categoria = [
        ("AR CONDICIONADOS", 2, 3, 4),
        ("BEBEDOUROS", 5, 6, None),
        (CATEGORIA_CHILLER, 7, 8, None),
    ]
    esquerda_por_categoria = {"AR CONDICIONADOS": 548640, "BEBEDOUROS": 4297680, CATEGORIA_CHILLER: 8046720}
    for categoria, idx_titulo, idx_lista, idx_overflow in blocos_categoria:
        pendentes = d["pendentes_por_categoria"].get(categoria, [])
        set_shape_text(shapes[idx_titulo], f"{categoria} ({len(pendentes)})")
        mostrados = pendentes[:MAX_ITENS_PENDENCIA_POR_CATEGORIA]
        resto = len(pendentes) - len(mostrados)
        set_shape_paragraphs(
            shapes[idx_lista],
            [f"{nome} \u2014 {detalhe}" for nome, detalhe in mostrados] or ["Nenhum pendente \u2014 categoria em dia."],
        )
        if idx_overflow is not None:
            overflow_shape = shapes[idx_overflow]
            if resto > 0:
                set_shape_text(overflow_shape, f"+ {resto} outro(s)")
            else:
                el = overflow_shape._element
                el.getparent().remove(el)

    # segunda passada: BEBEDOUROS/CHILLER não têm shape de overflow no
    # molde (só AR precisou dele em julho/2026) — se precisarem agora,
    # clona o shape de overflow da AR e reposiciona na coluna certa.
    shapes = list(slides[12].shapes)
    overflow_modelo_idx = next(
        (i for i, s in enumerate(shapes) if s.has_text_frame and s.text_frame.text.startswith("+ ")), None
    )
    if overflow_modelo_idx is not None:
        spTree = slides[12].shapes._spTree
        for categoria, _, _, idx_overflow in blocos_categoria:
            if idx_overflow is not None:
                continue
            pendentes = d["pendentes_por_categoria"].get(categoria, [])
            resto = len(pendentes) - MAX_ITENS_PENDENCIA_POR_CATEGORIA
            if resto > 0:
                modelo_el = list(slides[12].shapes)[overflow_modelo_idx]._element
                novo_el = copy.deepcopy(modelo_el)
                spTree.append(novo_el)
                xfrm = novo_el.find(f".//{qn('a:xfrm')}")
                off = xfrm.find(qn("a:off"))
                off.set("x", str(esquerda_por_categoria[categoria]))
                r = novo_el.find(f".//{qn('a:r')}")
                textos_t = r.findall(qn("a:t"))
                for extra in textos_t[1:]:
                    extra.getparent().remove(extra)
                textos_t[0].text = f"+ {resto} outro(s)"

    # --- slide 13: conclusão ---
    shape = list(slides[13].shapes)[1]
    paragrafos = shape.text_frame.paragraphs
    set_run_text(
        paragrafos[2],
        f"Em {d['nome_mes_ano']}, {d['preventivas_mes_qtd']} manuten\u00e7\u00f5es preventivas e "
        f"{d['corretivas_mes_qtd']} corretivas foram realizadas nos {d['total_contrato']} equipamentos "
        f"sob contrato, resultando em {d['pct_cumprimento']}% de cumprimento geral da periodicidade de PMOC.",
    )
    set_run_text(
        paragrafos[4],
        f"Pontos que exigem aten\u00e7\u00e3o da equipe t\u00e9cnica: {len(d['reincidencias'])} equipamento(s) com "
        f"reincid\u00eancia de falha no ano, {len(d['vencidos'])} bebedouro(s) com troca de filtro vencida, e "
        f"{d['total_pendente']} equipamento(s) ainda fora do prazo de preventiva \u2014 detalhados nas se\u00e7\u00f5es "
        f"anteriores.",
    )

    atualizar_rodapes(prs, d["nome_mes_ano"])

    prs.save(caminho_saida)


# ---------------------------------------------------------------------------
# 4) Upload — bucket relatorios-pmoc do Supabase (obrigatório, é de lá que
#    o painel /kpi-manutencao baixa) + Nextcloud (opcional, cópia de backup
#    na mesma pasta onde o processo antigo deixava o "latest.pptx").
# ---------------------------------------------------------------------------

NEXTCLOUD_PASTA_RELATORIO = {
    "Mangels": "REFRIGERAÇÃO AÇOS - RELATÓRIOS",
}

CONTENT_TYPE_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def enviar_para_supabase_storage(caminho_arquivo: str, nome_arquivo: str) -> None:
    headers = {
        "apikey": SUPABASE_SERVICE_ROLE_KEY,
        "Authorization": f"Bearer {SUPABASE_SERVICE_ROLE_KEY}",
        "Content-Type": CONTENT_TYPE_PPTX,
        "x-upsert": "true",
    }
    with open(caminho_arquivo, "rb") as f:
        resp = requests.post(
            f"{SUPABASE_URL}/storage/v1/object/relatorios-pmoc/{nome_arquivo}",
            headers=headers,
            data=f.read(),
        )
    resp.raise_for_status()


def enviar_para_nextcloud(caminho_arquivo: str, nome_arquivo: str, cliente: str) -> None:
    pasta = NEXTCLOUD_PASTA_RELATORIO.get(cliente)
    if not pasta:
        print(f"  (sem pasta de relatório no Nextcloud configurada pra {cliente} — pulando essa cópia)")
        return
    nextcloud_url = os.environ.get("NEXTCLOUD_URL")
    nextcloud_user = os.environ.get("NEXTCLOUD_USER")
    nextcloud_senha = os.environ.get("NEXTCLOUD_APP_PASSWORD")
    if not (nextcloud_url and nextcloud_user and nextcloud_senha):
        print("  (NEXTCLOUD_URL/USER/APP_PASSWORD não configurados — pulando cópia no Nextcloud)")
        return
    from webdav3.client import Client

    client = Client({
        "webdav_hostname": f"{nextcloud_url}/remote.php/dav/files/{nextcloud_user}",
        "webdav_login": nextcloud_user,
        "webdav_password": nextcloud_senha,
        "disable_check": True,
    })
    client.upload_sync(remote_path=f"{pasta}/{nome_arquivo}", local_path=caminho_arquivo)


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--cliente", required=True, choices=["Mangels", "Noel"])
    parser.add_argument("--mes", help="AAAA-MM (padrão: mês anterior ao atual)")
    parser.add_argument("--saida", help="Caminho do .pptx de saída (padrão: AAAA-MM.pptx nesta pasta)")
    parser.add_argument(
        "--upload", action="store_true",
        help="Depois de gerar, sobe pro Supabase Storage (e Nextcloud, se configurado).",
    )
    args = parser.parse_args()

    if args.mes:
        ano, mes = (int(x) for x in args.mes.split("-"))
    else:
        hoje = datetime.now()
        primeiro_dia_mes_atual = datetime(hoje.year, hoje.month, 1)
        mes_anterior = primeiro_dia_mes_atual - timedelta(days=1)
        ano, mes = mes_anterior.year, mes_anterior.month

    print(f"Gerando relat\u00f3rio PMOC de {args.cliente} \u2014 compet\u00eancia {mes:02d}/{ano}...")
    atendimentos = carregar_atendimentos(args.cliente)
    equipamentos = carregar_equipamentos(args.cliente)
    print(f"  {len(atendimentos)} atendimento(s), {len(equipamentos)} equipamento(s) carregado(s).")

    dados = calcular_relatorio(args.cliente, atendimentos, equipamentos, ano, mes)

    nome_arquivo = f"{ano:04d}-{mes:02d}.pptx"
    caminho_saida = args.saida or os.path.join(DIR_ATUAL, nome_arquivo)
    gerar_pptx(dados, caminho_saida)
    print(f"Relat\u00f3rio salvo em: {caminho_saida}")

    if args.upload:
        print(f"Enviando para o Supabase Storage (relatorios-pmoc/{nome_arquivo})...")
        enviar_para_supabase_storage(caminho_saida, nome_arquivo)
        print("  ok.")
        print("Enviando c\u00f3pia para o Nextcloud...")
        enviar_para_nextcloud(caminho_saida, nome_arquivo, args.cliente)
        print("  ok.")


if __name__ == "__main__":
    main()
